import os
import json
import subprocess
from collections import defaultdict
import argparse
from concurrent.futures import ThreadPoolExecutor
from docx import Document
import openpyxl
from pptx import Presentation
import PyPDF2
import logging
from typing import Dict, Optional

from langchain_huggingface import HuggingFaceEmbeddings

# Preload the embedding model once for efficiency using LangChain
_EMBED_MODEL = HuggingFaceEmbeddings(model_name="all-MiniLM-L6-v2")


def get_embedding(text: str) -> list[float]:
    """Return an embedding vector for the given text."""
    return _EMBED_MODEL.embed_query(text)


from .llm_provider import get_llm

# ------ CONFIG ------
# Configuration can be supplied via command line or environment variables.
# Environment variable fallbacks: FO_ROOT_DIR, FO_N_SAMPLE_FILES, FO_OLLAMA_MODEL


def parse_args():
    parser = argparse.ArgumentParser(description="Summarize folders using a local LLM")
    parser.add_argument(
        "--root",
        required=False,
        default=os.environ.get("FO_ROOT_DIR"),
        help="Root directory to analyze",
    )
    parser.add_argument(
        "--samples",
        type=int,
        default=int(os.environ.get("FO_N_SAMPLE_FILES", 10)),
        help="Number of sample files per folder",
    )
    parser.add_argument(
        "--model",
        default=os.environ.get("FO_OLLAMA_MODEL", "llama3"),
        help="Model name for the selected provider",
    )
    parser.add_argument(
        "--provider",
        choices=["ollama", "openai", "fireworks"],
        default=os.environ.get("FO_PROVIDER", "ollama"),
        help="LLM provider to use",
    )
    parser.add_argument(
        "--openai-api-key",
        dest="openai_api_key",
        default=os.environ.get("OPENAI_API_KEY"),
        help="API key for OpenAI provider",
    )
    parser.add_argument(
        "--fireworks-api-key",
        dest="fireworks_api_key",
        default=os.environ.get("FIREWORKS_API_KEY"),
        help="API key for Fireworks provider",
    )
    parser.add_argument(
        "--resume",
        action="store_true",
        help="Resume from existing folder_contexts.json",
    )
    parser.add_argument("--verbose", action="store_true", help="Enable verbose output")
    return parser.parse_args()


# --- File Extractors ---


def extract_text_txt(path, n_chars: Optional[int] = None):
    try:
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            text = f.read()
            return text if n_chars is None else text[:n_chars]
    except Exception:
        return ""


def extract_text_docx(path, n_chars: Optional[int] = None):
    try:
        doc = Document(path)
        text = "\n".join([p.text for p in doc.paragraphs])
        return text if n_chars is None else text[:n_chars]
    except Exception:
        return ""


def extract_text_pdf(path, n_chars: Optional[int] = None):
    try:
        reader = PyPDF2.PdfReader(path)
        text = ""
        for page in reader.pages:
            text += page.extract_text() or ""
        return text if n_chars is None else text[:n_chars]
    except Exception:
        return ""


def extract_text_xlsx(path, n_chars: Optional[int] = None):
    try:
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        text = ""
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                rowtext = "\t".join(
                    [str(cell) if cell is not None else "" for cell in row]
                )
                text += rowtext + "\n"
        return text if n_chars is None else text[:n_chars]
    except Exception:
        return ""


def extract_text_pptx(path, n_chars: Optional[int] = None):
    try:
        prs = Presentation(path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text if n_chars is None else text[:n_chars]
    except Exception:
        return ""


# Mapping of file extensions to their extraction functions
EXTRACTORS = {
    ".txt": extract_text_txt,
    ".md": extract_text_txt,
    ".csv": extract_text_txt,
    ".docx": extract_text_docx,
    ".pdf": extract_text_pdf,
    ".xlsx": extract_text_xlsx,
    ".xls": extract_text_xlsx,
    ".pptx": extract_text_pptx,
    ".ppt": extract_text_pptx,
}


def extract_text_file(path, n_chars: Optional[int] = None):
    ext = os.path.splitext(path)[1].lower()
    extractor = EXTRACTORS.get(ext)
    return extractor(path, n_chars) if extractor else ""


def get_sample_files(folder, n):
    files = [f for f in os.listdir(folder) if os.path.isfile(os.path.join(folder, f))]
    if not files:
        return []

    # Prefer recently modified files and distribute across extensions
    files.sort(key=lambda f: os.path.getmtime(os.path.join(folder, f)), reverse=True)
    by_ext = defaultdict(list)
    for f in files:
        by_ext[os.path.splitext(f)[1]].append(f)

    selected = []
    while len(selected) < n and by_ext:
        for ext in list(by_ext.keys()):
            if by_ext[ext]:
                selected.append(by_ext[ext].pop(0))
                if len(selected) >= n:
                    break
            if not by_ext[ext]:
                del by_ext[ext]
    return selected[:n]


def call_llm(prompt, llm, retries=1):
    for attempt in range(retries + 1):
        try:
            return llm.invoke(prompt).content.strip()
        except Exception as e:
            if attempt >= retries:
                logging.error("[llm error] %s", e)
                return ""
            logging.warning("LLM call failed, retrying...")


def get_file_summary(filepath, llm):
    content = extract_text_file(filepath, n_chars=2000)
    if not content:
        return "No extractable text."
    prompt = f"Summarize this file in one sentence for a folder classification system:\n{content}"
    summary = call_llm(prompt, llm=llm)
    return summary.strip()


# -- Main hierarchical logic --


def build_folder_tree(root_dir):
    """Map each folder to its children."""
    tree = defaultdict(list)
    all_folders = set()
    for folder, dirs, files in os.walk(root_dir):
        all_folders.add(folder)
        for d in dirs:
            child = os.path.join(folder, d)
            tree[folder].append(child)
    return tree, all_folders


def get_folders_in_bottom_up_order(tree, root_dir):
    """Return a list of folders, leaves first, root last."""
    visited = set()
    order = []

    def visit(folder):
        if folder in visited:
            return
        for child in tree[folder]:
            visit(child)
        visited.add(folder)
        order.append(folder)

    visit(root_dir)
    return order


def get_display_path(folder: str, root_dir: str) -> str:
    """Return folder path including root folder name for context."""
    folder = os.path.abspath(folder)
    root_dir = os.path.abspath(root_dir)
    root_name = os.path.basename(root_dir)
    if folder == root_dir:
        return root_name
    relative = os.path.relpath(folder, root_dir)
    return os.path.join(root_name, relative)


class FolderOrganizer:
    """High level organizer that manages folder summarization."""

    def __init__(
        self,
        root: str,
        *,
        samples: int = 10,
        model: str = "llama3",
        provider: str = "ollama",
        openai_api_key: Optional[str] = None,
        fireworks_api_key: Optional[str] = None,
        verbose: bool = False,
    ) -> None:
        self.root = root
        self.samples = samples
        self.model = model
        self.provider = provider
        self.openai_api_key = openai_api_key
        self.fireworks_api_key = fireworks_api_key
        logging.basicConfig(
            level=logging.DEBUG if verbose else logging.INFO, format="%(message)s"
        )
        self.llm = get_llm(
            model,
            provider,
            openai_api_key=openai_api_key,
            fireworks_api_key=fireworks_api_key,
        )
        self.tree: Dict[str, list] = {}
        self.order: list[str] = []
        self.out_json = os.path.join(self.root, "folder_contexts.json")
        self.out_vectors = os.path.join(self.root, "folder_vectors.json")

    def get_display_path(self, folder: str) -> str:
        return get_display_path(folder, self.root)

    def build_tree(self) -> Dict[str, list]:
        self.tree, _ = build_folder_tree(self.root)
        self.order = get_folders_in_bottom_up_order(self.tree, self.root)
        return self.tree

    def write_results(
        self, folder_contexts: Dict[str, str], folder_vectors: Dict[str, list]
    ) -> None:
        with open(self.out_json, "w", encoding="utf-8") as f:
            json.dump(folder_contexts, f, ensure_ascii=False, indent=2)
        with open(self.out_vectors, "w", encoding="utf-8") as f:
            json.dump(folder_vectors, f, ensure_ascii=False)

    def summarize_folders(self, *, resume: bool = False) -> None:
        if not self.tree:
            self.build_tree()

        if resume and os.path.exists(self.out_json):
            with open(self.out_json, "r", encoding="utf-8") as f:
                folder_contexts = json.load(f)
        else:
            folder_contexts = {}

        if resume and os.path.exists(self.out_vectors):
            with open(self.out_vectors, "r", encoding="utf-8") as f:
                folder_vectors: Dict[str, list] = json.load(f)
        else:
            folder_vectors = {}

        for folder in self.order:
            if folder in folder_contexts:
                logging.info("Skipping: %s", folder)
                continue

            logging.info("\nProcessing: %s", folder)
            folder_display = self.get_display_path(folder)
            sample_files = get_sample_files(folder, self.samples)
            sample_summaries = []

            files_to_process = sample_files[:3]
            with ThreadPoolExecutor() as ex:
                summaries = list(
                    ex.map(
                        lambda f: get_file_summary(os.path.join(folder, f), self.llm),
                        files_to_process,
                    )
                )
            for fname, summary in zip(files_to_process, summaries):
                sample_summaries.append(f"File: {fname}\nSummary: {summary}")

            embed_parts = [
                f"Path: {folder_display}",
                f"Folder: {os.path.basename(folder)}",
            ]
            all_files = [
                f for f in os.listdir(folder) if os.path.isfile(os.path.join(folder, f))
            ]
            for fname in all_files:
                snippet = extract_text_file(os.path.join(folder, fname))
                if snippet:
                    embed_parts.append(snippet)
            embed_text = "\n".join(embed_parts)
            folder_vectors[folder] = get_embedding(embed_text)

            child_contexts = []
            for child in self.tree[folder]:
                if child in folder_contexts:
                    child_display = self.get_display_path(child)
                    child_contexts.append(
                        f"Subfolder '{child_display}': {folder_contexts[child]}"
                    )

            prompt = f"""
You are an expert at understanding folder content and organization. The full
folder path provides important context. Subfolders inherit the meaning of the
entire path. For example, if the path contains  'Education' , treat every
subfolder as being about education . Your task is to
summarize the *main topic* of the folder below, ignoring any files that don't fit
the main theme.

Full folder path from the root: '{folder_display}'
Example file summaries (auto-generated):
{chr(10).join(sample_summaries)}
Subfolder context summaries:
{chr(10).join(child_contexts)}

Please summarize the main purpose or topic of this folder in 2-3 sentences,
taking into account its path-derived context, its own files, and the main themes
of its immediate subfolders (if any). If you see outliers, ignore them. Only
output the summary text.
"""
            logging.info("  Calling %s (%s)...", self.provider, self.model)
            folder_summary = call_llm(prompt, self.llm)
            logging.info("  => %s", folder_summary)

            folder_contexts[folder] = folder_summary
            self.write_results(folder_contexts, folder_vectors)

        logging.info("\nSaved folder contexts to %s", self.out_json)
        logging.info("Saved folder vectors to %s", self.out_vectors)


def main():
    args = parse_args()
    if not args.root:
        raise SystemExit("--root must be specified (or FO_ROOT_DIR env variable)")

    organizer = FolderOrganizer(
        root=args.root,
        samples=args.samples,
        model=args.model,
        provider=args.provider,
        openai_api_key=args.openai_api_key,
        fireworks_api_key=args.fireworks_api_key,
        verbose=args.verbose,
    )

    organizer.build_tree()
    organizer.summarize_folders(resume=args.resume)


if __name__ == "__main__":
    main()
