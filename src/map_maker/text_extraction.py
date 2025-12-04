from __future__ import annotations

import mimetypes
import xml.etree.ElementTree as ET
import zipfile
from pathlib import Path
from typing import Iterable, List, Tuple

from docx import Document
from pptx import Presentation
from pypdf import PdfReader

from .config import SAMPLE_BYTES

TEXT_EXTENSIONS = {
    ".txt",
    ".md",
    ".rtf",
    ".csv",
    ".json",
    ".yaml",
    ".yml",
    ".ini",
    ".log",
    ".py",
    ".java",
    ".js",
    ".ts",
    ".html",
    ".css",
    ".xml",
    ".xmind",
}

PDF_MIME = {"application/pdf"}
DOCX_MIME = {"application/vnd.openxmlformats-officedocument.wordprocessingml.document"}
PPTX_MIME = {"application/vnd.openxmlformats-officedocument.presentationml.presentation"}


def is_textual(path: Path) -> bool:
    if path.suffix.lower() in TEXT_EXTENSIONS:
        return True
    if path.suffix.lower() in {".pptx", ".ppt"}:
        return True
    mime, _ = mimetypes.guess_type(path)
    return bool(mime and (mime.startswith("text/") or mime in PDF_MIME or mime in DOCX_MIME or mime in PPTX_MIME))


def read_text_file(path: Path) -> str:
    try:
        return path.read_text(encoding="utf-8", errors="ignore")[:SAMPLE_BYTES]
    except UnicodeDecodeError:
        return path.read_text(encoding="latin-1", errors="ignore")[:SAMPLE_BYTES]


def read_pdf(path: Path) -> str:
    reader = PdfReader(path)
    pages = []
    for page in reader.pages[:3]:
        text = page.extract_text() or ""
        pages.append(text)
    return "\n".join(pages)[:SAMPLE_BYTES]


def read_docx(path: Path) -> str:
    document = Document(path)
    paragraphs = [p.text for p in document.paragraphs]
    return "\n".join(paragraphs)[:SAMPLE_BYTES]


def read_pptx(path: Path) -> str:
    """Extract text content from PowerPoint files"""
    presentation = Presentation(path)
    text_parts = []

    for slide in presentation.slides:
        # Extract title
        if slide.shapes.title:
            text_parts.append(slide.shapes.title.text)

        # Extract text from all text boxes
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text_parts.append(shape.text)

    return "\n".join(text_parts)[:SAMPLE_BYTES]


def read_xmind(path: Path) -> str:
    """Extract text content from XMind mind map files"""
    try:
        with zipfile.ZipFile(path, 'r') as zip_ref:
            # XMind files contain content.xml or content.json
            if 'content.xml' in zip_ref.namelist():
                content = zip_ref.read('content.xml').decode('utf-8', errors='ignore')
                # Parse XML and extract topic titles
                root = ET.fromstring(content)
                topics = []

                # Extract all topic titles (XMind uses 'topic' elements with 'title' attributes)
                for topic in root.iter():
                    if 'title' in topic.attrib:
                        topics.append(topic.attrib['title'])
                    # Also check for text content in elements
                    if topic.text and topic.text.strip():
                        topics.append(topic.text.strip())

                return "\n".join(topics)[:SAMPLE_BYTES]
            elif 'content.json' in zip_ref.namelist():
                # Newer XMind versions use JSON
                import json
                content = zip_ref.read('content.json').decode('utf-8', errors='ignore')
                data = json.loads(content)
                topics = []

                def extract_topics(node):
                    if isinstance(node, dict):
                        if 'title' in node:
                            topics.append(node['title'])
                        for value in node.values():
                            extract_topics(value)
                    elif isinstance(node, list):
                        for item in node:
                            extract_topics(item)

                extract_topics(data)
                return "\n".join(topics)[:SAMPLE_BYTES]
    except Exception:
        pass
    return ""


def safe_extract(path: Path) -> Tuple[str, List[str]]:
    errors: List[str] = []
    text = ""
    try:
        if path.suffix.lower() == ".pdf":
            text = read_pdf(path)
        elif path.suffix.lower() in {".docx", ".doc"}:
            text = read_docx(path)
        elif path.suffix.lower() in {".pptx", ".ppt"}:
            text = read_pptx(path)
        elif path.suffix.lower() == ".xmind":
            text = read_xmind(path)
        else:
            text = read_text_file(path)
    except Exception as exc:  # pylint: disable=broad-except
        errors.append(f"Failed to read {path.name}: {exc}")
    return text[:SAMPLE_BYTES], errors


def sample_files(files: Iterable[Path], limit: int) -> List[Path]:
    candidates = sorted(files, key=lambda p: p.stat().st_mtime)
    if len(candidates) <= limit:
        return candidates

    oldest = candidates[:3]
    newest = candidates[-3:]
    remaining_pool = [f for f in candidates if f not in oldest + newest]

    remaining_needed = max(0, limit - len(oldest) - len(newest))
    sampled: List[Path] = []
    if remaining_needed > 0 and remaining_pool:
        stride = max(1, len(remaining_pool) // remaining_needed)
        for idx in range(remaining_needed):
            pick_index = min(idx * stride, len(remaining_pool) - 1)
            sampled.append(remaining_pool[pick_index])

    combined = list(dict.fromkeys(oldest + sampled + newest))
    return combined[:limit]
