import os
import json
import subprocess
from collections import defaultdict, Counter
from docx import Document
import openpyxl
from pptx import Presentation
import PyPDF2

# ------ CONFIG ------
ROOT_DIR = "/path/to/your/root/folder"   # <--- Change this!
N_SAMPLE_FILES = 10
OLLAMA_MODEL = "llama3"                  # or another local model name

# --- File Extractors ---

def extract_text_txt(path, n_chars=2000):
    try:
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            return f.read(n_chars)
    except Exception:
        return ""

def extract_text_docx(path, n_chars=2000):
    try:
        doc = Document(path)
        text = "\n".join([p.text for p in doc.paragraphs])
        return text[:n_chars]
    except Exception:
        return ""

def extract_text_pdf(path, n_chars=2000):
    try:
        reader = PyPDF2.PdfReader(path)
        text = ""
        for page in reader.pages[:5]:
            text += page.extract_text() or ""
        return text[:n_chars]
    except Exception:
        return ""

def extract_text_xlsx(path, n_chars=2000):
    try:
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        text = ""
        for ws in wb.worksheets[:2]:
            for row in ws.iter_rows(min_row=1, max_row=15, values_only=True):
                rowtext = "\t".join([str(cell) if cell is not None else "" for cell in row])
                text += rowtext + "\n"
        return text[:n_chars]
    except Exception:
        return ""

def extract_text_pptx(path, n_chars=2000):
    try:
        prs = Presentation(path)
        text = ""
        for i, slide in enumerate(prs.slides):
            if i > 5: break
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text[:n_chars]
    except Exception:
        return ""

def extract_text_file(path, n_chars=2000):
    ext = os.path.splitext(path)[1].lower()
    if ext in [".txt", ".md", ".csv"]:
        return extract_text_txt(path, n_chars)
    elif ext == ".docx":
        return extract_text_docx(path, n_chars)
    elif ext == ".pdf":
        return extract_text_pdf(path, n_chars)
    elif ext in [".xlsx", ".xls"]:
        return extract_text_xlsx(path, n_chars)
    elif ext in [".pptx", ".ppt"]:
        return extract_text_pptx(path, n_chars)
    else:
        return ""

def get_sample_files(folder, n=N_SAMPLE_FILES):
    files = [f for f in os.listdir(folder) if os.path.isfile(os.path.join(folder, f))]
    if not files:
        return []
    exts = [os.path.splitext(f)[1] for f in files]
    most_common_ext = Counter(exts).most_common(1)[0][0] if exts else ""
    rep_files = [f for f in files if os.path.splitext(f)[1] == most_common_ext]
    if len(rep_files) < n:
        rep_files = files
    return rep_files[:n]

def call_ollama(prompt, model=OLLAMA_MODEL):
    proc = subprocess.run(
        ["ollama", "run", model],
        input=prompt.encode("utf-8"),
        stdout=subprocess.PIPE
    )
    return proc.stdout.decode("utf-8").strip()

def get_file_summary(filepath, ollama_model=OLLAMA_MODEL):
    content = extract_text_file(filepath)
    if not content:
        return "No extractable text."
    prompt = f"Summarize this file in one sentence for a folder classification system:\n{content}"
    summary = call_ollama(prompt, model=ollama_model)
    return summary.strip()

# -- Main hierarchical logic --

def build_folder_tree(root_dir):
    """Map each folder to its children."""
    tree = defaultdict(list)
    all_folders = set()
    for folder, dirs, files in os.walk(root_dir):
        all_folders.add(folder)
        for d in dirs:
            child = os.path.join(folder, d)
            tree[folder].append(child)
    return tree, all_folders

def get_folders_in_bottom_up_order(tree, root_dir):
    """Return a list of folders, leaves first, root last."""
    visited = set()
    order = []

    def visit(folder):
        if folder in visited:
            return
        for child in tree[folder]:
            visit(child)
        visited.add(folder)
        order.append(folder)

    visit(root_dir)
    return order

def main():
    # Build folder hierarchy
    tree, all_folders = build_folder_tree(ROOT_DIR)
    process_order = get_folders_in_bottom_up_order(tree, ROOT_DIR)

    folder_contexts = {}
    for folder in process_order:
        print(f"\nProcessing: {folder}")
        # --- File-based summaries
        sample_files = get_sample_files(folder)
        sample_summaries = []
        for fname in sample_files[:3]:  # For prompt size, use up to 3 samples
            full_path = os.path.join(folder, fname)
            summary = get_file_summary(full_path)
            sample_summaries.append(f"File: {fname}\nSummary: {summary}")

        # --- Child context summaries
        child_contexts = []
        for child in tree[folder]:
            if child in folder_contexts:
                child_contexts.append(f"Subfolder '{os.path.basename(child)}': {folder_contexts[child]}")

        # --- Build context prompt
        prompt = f"""
You are an expert at understanding folder content and organization. Your task is to summarize the *main topic* of the folder below, ignoring any files that don't fit the main theme.

Folder name: '{os.path.basename(folder)}'
Example file summaries (auto-generated): 
{chr(10).join(sample_summaries)}
Subfolder context summaries:
{chr(10).join(child_contexts)}

Please summarize the main purpose or topic of this folder in 2-3 sentences, taking into account both its own files and the main themes of its immediate subfolders (if any). If you see outliers, ignore them. Only output the summary text.
"""
        print(f"  Calling Ollama ({OLLAMA_MODEL})...")
        folder_summary = call_ollama(prompt)
        print(f"  => {folder_summary}")

        folder_contexts[folder] = folder_summary

    out_json = os.path.join(ROOT_DIR, "folder_contexts.json")
    with open(out_json, "w", encoding="utf-8") as f:
        json.dump(folder_contexts, f, ensure_ascii=False, indent=2)
    print(f"\nSaved folder contexts to {out_json}")

if __name__ == "__main__":
    main()
